import os
import json
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from googleapiclient.http import MediaFileUpload

from src.services.google_auth import get_google_services
from src.services.ai_generator import generate_ai_content

# PDF and Word document text extraction libraries
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

def extract_text_from_file(file_path: str) -> str:
    """
    Extracts text content from a local PDF, DOCX, or TXT file.
    """
    if not os.path.exists(file_path):
        return ""
    
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        if ext == ".pdf":
            if pypdf:
                reader = pypdf.PdfReader(file_path)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
        elif ext in [".docx", ".doc"]:
            if docx:
                doc = docx.Document(file_path)
                for p in doc.paragraphs:
                    if p.text:
                        text += p.text + "\n"
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception as e:
        print(f"[SlidesService] Error extracting text from {file_path}: {e}")
        
    return text.strip()

def analyze_document_topics(file_path: str) -> dict:
    """
    Analyzes an uploaded document and extracts table of contents / chapter list / main topics
    so Zeyra can ask the user for specific guidelines on which chapters/topics to convert to slides.
    """
    text = extract_text_from_file(file_path)
    if not text:
        return {"document_name": os.path.basename(file_path), "topics": []}
    
    snippet = text[:5000]
    prompt = f"""
Analyze this document snippet and extract the main chapters, sections, or topics present:
Document Filename: {os.path.basename(file_path)}
Content Snippet:
{snippet}

Return JSON with schema:
{{
  "document_title": "Title or main subject of the document",
  "chapters_or_topics": [
    "Chapter 1 / Topic 1 Name",
    "Chapter 2 / Topic 2 Name",
    "Chapter 3 / Topic 3 Name"
  ]
}}
"""
    try:
        raw_json = generate_ai_content(prompt, response_mime_type="application/json")
        clean_str = raw_json.strip().strip("`").replace("json\n", "")
        return json.loads(clean_str)
    except Exception as e:
        print(f"[SlidesService] Error analyzing document topics: {e}")
        return {
            "document_title": os.path.basename(file_path),
            "chapters_or_topics": ["General Overview & Main Topics"]
        }

def generate_slides_outline(input_text_or_topic: str) -> dict:
    """
    Uses Gemini AI to structure raw lecture text, document content, or specific chapter into a structured presentation outline JSON.
    """
    prompt = f"""
You are an expert academic curriculum designer and executive presentation creator.
Create a high-impact, educational lecture presentation outline based strictly on the following input material / chapter guidelines:

Input Material / Chapter / Guidelines:
"{input_text_or_topic[:6000]}"

Structure the output as strict JSON with this exact schema:
{{
  "presentation_title": "Clean, Catchy Presentation Title",
  "subtitle": "Lecture Subtitle / Subject / Chapter Name",
  "instructor": "Muhammad Hamza",
  "objectives": [
    "Learning objective 1",
    "Learning objective 2",
    "Learning objective 3"
  ],
  "slides": [
    {{
      "title": "Slide Title 1",
      "bullet_points": [
        "Key concept 1 with clear explanation",
        "Key concept 2 with clear explanation",
        "Key concept 3 with clear explanation"
      ],
      "takeaway": "Core takeaway or student note"
    }},
    {{
      "title": "Slide Title 2",
      "bullet_points": [
        "Key point 1",
        "Key point 2",
        "Key point 3"
      ],
      "takeaway": "Core takeaway summary"
    }}
  ],
  "summary_points": [
    "Summary takeaway 1",
    "Summary takeaway 2"
  ],
  "discussion_question": "A thought-provoking question for student Q&A discussion"
}}

Provide 4 to 7 content slides. Ensure content is educational, professional, and readable. Return ONLY valid JSON.
"""
    raw_json = generate_ai_content(prompt, response_mime_type="application/json")
    try:
        clean_str = raw_json.strip()
        if clean_str.startswith("```json"):
            clean_str = clean_str[7:]
        if clean_str.startswith("```"):
            clean_str = clean_str[3:]
        if clean_str.endswith("```"):
            clean_str = clean_str[:-3]
        return json.loads(clean_str.strip())
    except Exception as e:
        print(f"[SlidesService] Error parsing Gemini outline JSON: {e}")
        return {
            "presentation_title": "Academic Lecture Presentation",
            "subtitle": "Generated Lecture Notes",
            "instructor": "Muhammad Hamza",
            "objectives": ["Understand core principles", "Apply concepts to practical examples"],
            "slides": [
                {
                    "title": "Introduction to Topic",
                    "bullet_points": ["Overview of fundamental concepts", "Key terms and definitions", "Real-world significance"],
                    "takeaway": "Grasp the core framework before moving into implementation."
                }
            ],
            "summary_points": ["Review main takeaways", "Prepare for next class"],
            "discussion_question": "What are the practical applications of today's lecture?"
        }

def create_styled_pptx(outline: dict, output_path: str = "generated_lecture.pptx") -> str:
    """
    Creates a styled PPTX presentation using python-pptx with a Modern High-Contrast Light Snow / Pure White Theme
    specifically designed for high visibility from a distance in classroom/lecture halls.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Color Palette: Modern High-Contrast Light / White Executive Theme
    COLOR_BG = RGBColor(255, 255, 255)       # #FFFFFF Soft Pure White Background
    COLOR_CARD = RGBColor(241, 245, 249)     # #F1F5F9 Slate Light Ice Card
    COLOR_TITLE = RGBColor(15, 23, 42)       # #0F172A High-Contrast Deep Midnight Navy Header
    COLOR_BODY = RGBColor(30, 41, 59)        # #1E293B Dark Charcoal Text (Easy distance reading)
    COLOR_ACCENT = RGBColor(37, 99, 235)     # #2563EB Royal Blue Accent
    COLOR_MUTED = RGBColor(100, 116, 139)    # #64748B Slate Muted Text

    blank_layout = prs.slide_layouts[6]

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

    # 1. Slide 1: Title Slide (Clean White & High Contrast Dark Header)
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)

    # Accent Top Bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.5), Inches(2.5), Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_ACCENT
    top_bar.line.fill.background()

    # Title box
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = outline.get("presentation_title", "Lecture Presentation")
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = outline.get("subtitle", "")
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_ACCENT
    p2.space_before = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = f"Instructor: {outline.get('instructor', 'Muhammad Hamza')}"
    p3.font.size = Pt(18)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(28)

    # 2. Slide 2: Learning Objectives
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)

    tb2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.7), Inches(11.333), Inches(1.0))
    tf2 = tb2.text_frame
    p = tf2.paragraphs[0]
    p.text = "🎯 Learning Objectives"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE

    objs = outline.get("objectives", [])
    for idx, obj in enumerate(objs):
        top_pos = 1.9 + (idx * 1.3)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(top_pos), Inches(11.333), Inches(1.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.fill.background()

        ctb = slide2.shapes.add_textbox(Inches(1.2), Inches(top_pos + 0.15), Inches(10.9), Inches(0.7))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = f"•  {obj}"
        cp.font.size = Pt(20)
        cp.font.bold = True
        cp.font.color.rgb = COLOR_BODY

    # 3. Topic Content Slides
    for slide_data in outline.get("slides", []):
        cslide = prs.slides.add_slide(blank_layout)
        set_slide_bg(cslide)

        # Header Title
        htb = cslide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
        htf = htb.text_frame
        hp = htf.paragraphs[0]
        hp.text = slide_data.get("title", "Topic Overview")
        hp.font.size = Pt(30)
        hp.font.bold = True
        hp.font.color.rgb = COLOR_TITLE

        # Content Card
        bcard = cslide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.7), Inches(11.333), Inches(4.3))
        bcard.fill.solid()
        bcard.fill.fore_color.rgb = COLOR_CARD
        bcard.line.fill.background()

        btb = cslide.shapes.add_textbox(Inches(1.3), Inches(1.9), Inches(10.7), Inches(3.9))
        btf = btb.text_frame
        btf.word_wrap = True

        bullets = slide_data.get("bullet_points", [])
        for b_idx, bullet in enumerate(bullets):
            bp = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
            bp.text = f"•  {bullet}"
            bp.font.size = Pt(20)
            bp.font.color.rgb = COLOR_BODY
            if b_idx > 0:
                bp.space_before = Pt(14)

        # Takeaway footer line
        takeaway = slide_data.get("takeaway")
        if takeaway:
            tbox = cslide.shapes.add_textbox(Inches(1.0), Inches(6.3), Inches(11.333), Inches(0.8))
            ttf = tbox.text_frame
            tp = ttf.paragraphs[0]
            tp.text = f"💡 Key Takeaway: {takeaway}"
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = COLOR_ACCENT

    # 4. Slide Final: Summary & Discussion
    fslide = prs.slides.add_slide(blank_layout)
    set_slide_bg(fslide)

    ftb = fslide.shapes.add_textbox(Inches(1.0), Inches(0.7), Inches(11.333), Inches(1.0))
    ftf = ftb.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "📌 Summary & Student Discussion"
    fp.font.size = Pt(32)
    fp.font.bold = True
    fp.font.color.rgb = COLOR_TITLE

    # Summary box
    scard = fslide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.2))
    scard.fill.solid()
    scard.fill.fore_color.rgb = COLOR_CARD
    scard.line.fill.background()

    stb = fslide.shapes.add_textbox(Inches(1.3), Inches(2.0), Inches(10.7), Inches(2.8))
    stf = stb.text_frame
    stf.word_wrap = True

    sum_points = outline.get("summary_points", [])
    for idx, sp_text in enumerate(sum_points):
        sp = stf.paragraphs[0] if idx == 0 else stf.add_paragraph()
        sp.text = f"•  {sp_text}"
        sp.font.size = Pt(20)
        sp.font.color.rgb = COLOR_BODY
        if idx > 0:
            sp.space_before = Pt(12)

    # Q&A Discussion Prompt
    q_prompt = outline.get("discussion_question")
    if q_prompt:
        qbox = fslide.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(11.333), Inches(1.5))
        qtf = qbox.text_frame
        qp = qtf.paragraphs[0]
        qp.text = f"❓ Q&A Discussion Prompt: {q_prompt}"
        qp.font.size = Pt(20)
        qp.font.bold = True
        qp.font.color.rgb = COLOR_ACCENT

    prs.save(output_path)
    return output_path

def create_google_slides_presentation(input_content_or_filepath: str) -> dict:
    """
    Complete pipeline: parses input (text or uploaded PDF/DOCX/TXT file),
    generates AI outline, builds PPTX presentation with High-Contrast White Theme,
    uploads & converts to native Google Slides, sets public edit permissions,
    and returns a clean, direct clickable Google Slides link.
    """
    # 1. Check if input is a local file path
    input_text = ""
    is_file = os.path.exists(str(input_content_or_filepath))
    
    if is_file:
        print(f"[SlidesService] Extracting text from file: {input_content_or_filepath}")
        input_text = extract_text_from_file(input_content_or_filepath)
        if not input_text:
            input_text = f"Lecture on {os.path.basename(input_content_or_filepath)}"
    else:
        input_text = str(input_content_or_filepath)

    print("[SlidesService] Generating AI presentation outline...")
    outline = generate_slides_outline(input_text)
    pres_title = outline.get("presentation_title", "Academic Lecture Presentation")
    
    # 2. Build local PPTX file
    sanitized_filename = f"Lecture_{pres_title.replace(' ', '_')[:30]}.pptx"
    pptx_path = os.path.join(os.getcwd(), sanitized_filename)
    create_styled_pptx(outline, output_path=pptx_path)
    print(f"[SlidesService] Local PPTX created: {pptx_path}")

    # 3. Upload & convert to Google Slides via Google Drive API
    sheets_service, gmail_service, docs_service, drive_service = get_google_services()
    if not drive_service:
        return {
            "success": False,
            "error": "Google Drive authentication unavailable. Please re-authenticate."
        }

    try:
        media = MediaFileUpload(
            pptx_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            resumable=True
        )
        file_metadata = {
            "name": f"🎓 {pres_title}",
            "mimeType": "application/vnd.google-apps.presentation"  # Converts PPTX directly to native Google Slides
        }
        gfile = drive_service.files().create(
            body=file_metadata,
            media_body=media,
            fields="id, name, webViewLink"
        ).execute()

        f_id = gfile.get("id")

        # Set public/shareable edit permissions for anyone with the link
        try:
            drive_service.permissions().create(
                fileId=f_id,
                body={"type": "anyone", "role": "writer"}
            ).execute()
        except Exception as pe:
            print(f"[SlidesService] Permission warning: {pe}")

        # Construct clean, un-mangled direct Google Slides edit link
        slides_direct_url = f"https://docs.google.com/presentation/d/{f_id}/edit?usp=sharing"
        slide_count = len(outline.get("slides", [])) + 3

        return {
            "success": True,
            "title": pres_title,
            "slides_count": slide_count,
            "file_id": f_id,
            "url": slides_direct_url,
            "pptx_file": pptx_path,
            "outline": outline
        }
    except Exception as e:
        print(f"[SlidesService] Error uploading presentation to Google Drive: {e}")
        return {
            "success": False,
            "error": str(e)
        }
